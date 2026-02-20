#!/usr/bin/env python3
"""スライド生成スクリプト - JSONからPPTXを生成する

使い方:
    python3 slide_generator.py <JSONファイルパス>

JSONスキーマ:
    {
      "title": "プレゼンタイトル",
      "output_dir": "slides",
      "style": "default",
      "slides": [
        {"layout": "title", "title": "...", "subtitle": "..."},
        {"layout": "image_caption", "title": "...", "body": "...",
         "bullets": ["...", "..."], "image_path": "/path/to/image.png"}
      ]
    }
"""

import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def load_style(style_name: str = "default") -> dict:
    style_path = Path(__file__).parent / "slide_style.json"
    with open(style_path, encoding="utf-8") as f:
        styles = json.load(f)
    return styles.get(style_name, styles["default"])


def hex_to_rgb(hex_str: str) -> RGBColor:
    """'2563EB' -> RGBColor(0x25, 0x63, 0xEB)"""
    return RGBColor(
        int(hex_str[0:2], 16),
        int(hex_str[2:4], 16),
        int(hex_str[4:6], 16),
    )


def _add_text(slide, left, top, width, height, text, font_size, bold=False, color_hex=None, align=PP_ALIGN.LEFT, word_wrap=True):
    """テキストボックスを追加するヘルパー"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color_hex:
        run.font.color.rgb = hex_to_rgb(color_hex)
    else:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return txBox, tf


def _set_background(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def create_title_slide(prs: Presentation, data: dict, style: dict):
    """タイトルスライド: フル背景色 + 中央テキスト"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_background(slide, style["primary_color"])

    w = prs.slide_width
    h = prs.slide_height

    _add_text(
        slide,
        Inches(1), h * 0.3,
        w - Inches(2), Inches(2),
        data.get("title", ""),
        style["title_font_size"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_text(
            slide,
            Inches(1), h * 0.55,
            w - Inches(2), Inches(1.5),
            subtitle,
            style["body_font_size"] + 4,
            align=PP_ALIGN.CENTER,
        )


def create_image_caption_slide(prs: Presentation, data: dict, style: dict):
    """画像 + テキストスライド: 上部ヘッダー + 左テキスト / 右画像"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_background(slide, style["background_color"])

    w = prs.slide_width
    h = prs.slide_height
    header_h = Inches(1.0)

    # 上部ヘッダーバー
    header = slide.shapes.add_shape(1, 0, 0, w, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(style["primary_color"])
    header.line.fill.background()

    # ヘッダー内タイトル
    _add_text(
        slide,
        Inches(0.3), Inches(0.1),
        w - Inches(0.6), header_h - Inches(0.2),
        data.get("title", ""),
        24,
        bold=True,
        align=PP_ALIGN.LEFT,
    )

    content_top = header_h + Inches(0.25)
    content_h = h - header_h - Inches(0.4)

    image_path = data.get("image_path")
    has_image = image_path and Path(image_path).exists()

    body = data.get("body", "")
    bullets = data.get("bullets", [])

    if has_image:
        text_w = w * 0.38
        img_left = w * 0.40
        img_w = w * 0.58

        # テキスト（左カラム）
        if body or bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.3), content_top,
                text_w - Inches(0.2), content_h,
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            lines = ([body] if body else []) + [f"• {b}" for b in bullets]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(8)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(style["body_font_size"])
                run.font.color.rgb = hex_to_rgb(style["text_color"])

        # 画像（右カラム）
        slide.shapes.add_picture(image_path, img_left, content_top, img_w, content_h)

    else:
        # 画像なし: テキスト全幅
        if body or bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.5), content_top,
                w - Inches(1.0), content_h,
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            lines = ([body] if body else []) + [f"• {b}" for b in bullets]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(8)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(style["body_font_size"])
                run.font.color.rgb = hex_to_rgb(style["text_color"])


LAYOUT_HANDLERS = {
    "title": create_title_slide,
    "image_caption": create_image_caption_slide,
    "bullet_points": create_image_caption_slide,
    "content": create_image_caption_slide,
}


def generate_pptx(json_path: str) -> Path:
    with open(json_path, encoding="utf-8") as f:
        data = json.load(f)

    style = load_style(data.get("style", "default"))

    prs = Presentation()
    prs.slide_width = Inches(style["slide_width_inches"])
    prs.slide_height = Inches(style["slide_height_inches"])

    for slide_data in data.get("slides", []):
        layout = slide_data.get("layout", "image_caption")
        handler = LAYOUT_HANDLERS.get(layout, create_image_caption_slide)
        handler(prs, slide_data, style)

    # 出力パス決定
    output_dir = Path(data.get("output_dir", "slides"))
    output_dir.mkdir(parents=True, exist_ok=True)

    title = data.get("title", "slides")
    safe_title = "".join(c for c in title if c.isalnum() or c in " \-_").strip()
    date_str = datetime.now().strftime("%Y%m%d")
    output_path = output_dir / f"{date_str}-{safe_title}.pptx"

    if output_path.exists():
        i = 2
        while output_path.exists():
            output_path = output_dir / f"{date_str}-{safe_title}-{i}.pptx"
            i += 1

    prs.save(output_path)
    print(f"保存完了: {output_path}")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("使い方: python3 slide_generator.py <JSONファイルパス>", file=sys.stderr)
        sys.exit(1)
    generate_pptx(sys.argv[1])
