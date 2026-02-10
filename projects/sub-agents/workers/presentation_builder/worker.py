#!/usr/bin/env python3
"""
Presentation Builder Worker - プレゼン資料作成ワーカー

役割: PPTXファイルの生成、スライド構成の提案、コンテンツの自動配置
対象: コンサルタント業務での提案書・プレゼン資料作成
"""

import json
import os
from typing import Dict, List, Any, Optional

# python-pptxが利用可能かチェック
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx がインストールされていません")
    print("   インストール: pip install python-pptx")


class PresentationBuilderWorker:
    """プレゼン資料作成ワーカー"""

    def __init__(self):
        """初期化"""
        self.name = "presentation_builder"
        self.description = "PPTX資料の作成を担当"
        self.skills = ["presentation-design", "visual-design", "storytelling"]
        self.layouts = self._load_layouts()

    def _load_layouts(self) -> Dict[str, Any]:
        """レイアウトパターンをロード"""
        layouts_path = os.path.join(os.path.dirname(__file__), "layouts", "layouts.json")

        default_layouts = {
            "title": {
                "name": "タイトルスライド",
                "description": "プレゼンの表紙",
                "elements": ["title", "subtitle", "author", "date"]
            },
            "section": {
                "name": "セクション区切り",
                "description": "大きな区切りを示すスライド",
                "elements": ["section_title", "section_number"]
            },
            "content": {
                "name": "コンテンツスライド",
                "description": "通常のコンテンツ",
                "elements": ["title", "body"]
            },
            "two_column": {
                "name": "2カラム",
                "description": "左右に内容を配置",
                "elements": ["title", "left_content", "right_content"]
            },
            "bullet_points": {
                "name": "箇条書き",
                "description": "ポイントを列挙",
                "elements": ["title", "bullets"]
            },
            "image_caption": {
                "name": "画像+説明",
                "description": "ビジュアル重視",
                "elements": ["title", "image", "caption"]
            },
            "conclusion": {
                "name": "まとめ",
                "description": "結論・まとめスライド",
                "elements": ["title", "key_points", "next_steps"]
            }
        }

        try:
            with open(layouts_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            return default_layouts

    def execute(self, task: Dict[str, Any]) -> Dict[str, Any]:
        """タスクを実行"""
        print(f"\n[*] {self.name}: タスクを実行中...")
        print(f"   タスク: {task.get('description', 'N/A')}")

        # タスクの種類を判定
        task_type = task.get("type", "create_presentation")

        if task_type == "create_presentation":
            return self._create_presentation(task)
        elif task_type == "suggest_structure":
            return self._suggest_structure(task)
        elif task_type == "generate_slide":
            return self._generate_slide(task)
        else:
            return {
                "status": "error",
                "worker": self.name,
                "error": f"不明なタスクタイプ: {task_type}"
            }

    def _suggest_structure(self, task: Dict[str, Any]) -> Dict[str, Any]:
        """スライド構成を提案"""
        topic = task.get("topic", "プレゼンテーション")
        audience = task.get("audience", "一般")
        duration = task.get("duration", 15)  # 分

        # スライド数の目安（1分あたり1スライド）
        estimated_slides = duration

        structure = [
            {
                "slide_number": 1,
                "layout": "title",
                "title": topic,
                "description": "表紙スライド",
                "notes": "会社名、発表者名、日付を含める"
            },
            {
                "slide_number": 2,
                "layout": "content",
                "title": "アジェンダ",
                "description": "本日の流れ",
                "notes": "3-5個のトピックを列挙"
            },
            {
                "slide_number": 3,
                "layout": "section",
                "title": "背景・課題",
                "description": "現状の問題点",
                "notes": "具体的な数字やデータで裏付ける"
            }
        ]

        # メインコンテンツ（スライド数の60%）
        main_slides = int(estimated_slides * 0.6)
        for i in range(main_slides):
            structure.append({
                "slide_number": len(structure) + 1,
                "layout": "content",
                "title": f"ポイント {i + 1}",
                "description": "主要な内容",
                "notes": "1スライド1メッセージの原則"
            })

        # まとめ
        structure.extend([
            {
                "slide_number": len(structure) + 1,
                "layout": "conclusion",
                "title": "まとめ",
                "description": "重要ポイントの再確認",
                "notes": "3つの要点に絞る"
            },
            {
                "slide_number": len(structure) + 1,
                "layout": "content",
                "title": "次のアクション",
                "description": "具体的な行動計画",
                "notes": "期限と担当を明記"
            }
        ])

        result = {
            "status": "success",
            "worker": self.name,
            "output": {
                "topic": topic,
                "audience": audience,
                "duration": duration,
                "estimated_slides": len(structure),
                "structure": structure,
                "recommendations": [
                    "1スライド1メッセージを心がける",
                    "ビジュアル（図・グラフ）を積極的に使用",
                    "文字数は最小限に（1スライド40文字以内が理想）",
                    "聴衆に合わせた専門用語の使用レベルを調整"
                ]
            },
            "logs": [
                f"トピック: {topic}",
                f"想定時間: {duration}分",
                f"提案スライド数: {len(structure)}枚",
                "構成を生成しました"
            ]
        }

        print(f"   [+] スライド構成を提案しました（{len(structure)}枚）")
        return result

    def _create_presentation(self, task: Dict[str, Any]) -> Dict[str, Any]:
        """PPTXファイルを作成"""
        if not PPTX_AVAILABLE:
            return {
                "status": "error",
                "worker": self.name,
                "error": "python-pptx がインストールされていません",
                "suggestion": "pip install python-pptx を実行してください"
            }

        topic = task.get("topic", "プレゼンテーション")
        slides_content = task.get("slides", [])
        output_path = task.get("output_path", f"{topic}.pptx")

        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # スライドを追加
            for slide_data in slides_content:
                self._add_slide(prs, slide_data)

            # 保存
            prs.save(output_path)

            result = {
                "status": "success",
                "worker": self.name,
                "output": {
                    "file_path": output_path,
                    "slide_count": len(slides_content),
                    "file_size": os.path.getsize(output_path) if os.path.exists(output_path) else 0
                },
                "logs": [
                    f"PPTXファイルを作成しました: {output_path}",
                    f"スライド数: {len(slides_content)}"
                ]
            }

            print(f"   [+] PPTX作成完了: {output_path}")
            return result

        except Exception as e:
            return {
                "status": "error",
                "worker": self.name,
                "error": str(e)
            }

    def _add_slide(self, prs: Any, slide_data: Dict[str, Any]):
        """スライドを追加"""
        layout_type = slide_data.get("layout", "content")

        # 空白レイアウトを使用
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # タイトルを追加
        if "title" in slide_data:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True

        # コンテンツを追加
        if "content" in slide_data:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(9), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.text = slide_data["content"]
            content_frame.paragraphs[0].font.size = Pt(18)

        # 箇条書きを追加
        if "bullets" in slide_data:
            bullets_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(9), Inches(5)
            )
            bullets_frame = bullets_box.text_frame
            for bullet in slide_data["bullets"]:
                p = bullets_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)

    def _generate_slide(self, task: Dict[str, Any]) -> Dict[str, Any]:
        """1枚のスライドコンテンツを生成"""
        slide_type = task.get("slide_type", "content")
        topic = task.get("topic", "")

        layout = self.layouts.get(slide_type, self.layouts["content"])

        slide_content = {
            "layout": slide_type,
            "title": topic,
            "elements": layout["elements"],
            "recommendations": [
                f"レイアウト: {layout['name']}",
                f"用途: {layout['description']}"
            ]
        }

        result = {
            "status": "success",
            "worker": self.name,
            "output": slide_content,
            "logs": [f"スライドコンテンツを生成: {slide_type}"]
        }

        return result


def main():
    """テスト用メイン関数"""
    worker = PresentationBuilderWorker()

    # テスト1: スライド構成の提案
    print("\n" + "="*60)
    print("テスト1: スライド構成の提案")
    print("="*60)

    task1 = {
        "type": "suggest_structure",
        "topic": "新規事業提案",
        "audience": "経営陣",
        "duration": 20
    }

    result1 = worker.execute(task1)

    if result1["status"] == "success":
        print(f"\n提案スライド数: {result1['output']['estimated_slides']}")
        print("\n構成:")
        for slide in result1['output']['structure'][:5]:  # 最初の5枚を表示
            print(f"  {slide['slide_number']}. {slide['title']} ({slide['layout']})")
        print("  ...")

        print("\n推奨事項:")
        for rec in result1['output']['recommendations']:
            print(f"  - {rec}")


if __name__ == "__main__":
    main()
